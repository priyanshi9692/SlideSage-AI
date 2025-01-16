import pandas as pd
from pptx import Presentation

class PPTDataPreprocessing:

    def __init__(self, file_path):
        self.file_path = file_path
        self.presentation = Presentation(file_path)

    def extract_tables_from_slide(self, slide):
        """
        Extracts tables from a PowerPoint slide and returns them as a list of pandas DataFrames.

        Args:
            slide (pptx.slide.Slide): A slide object from the PowerPoint presentation.

        Returns:
            list: A list of pandas DataFrames representing the extracted tables, or an empty list if no tables are found.
        """
        tables = []
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                data = []
                for row in table.rows:
                    data.append([cell.text for cell in row.cells])
                if data:
                    df = pd.DataFrame(data[1:], columns=data[0])  # Assuming first row is header
                    tables.append(df)
        return tables

    def extract_content_from_ppt(self):
        """
        Extracts all textual and tabular content from the presentation, combining slide text and tables into a formatted string.

        Returns:
            str: A single string containing the content of all slides, including text and tables.
        """
        overall_content = []
        for idx, slide in enumerate(self.presentation.slides):
            slide_content = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_content.append(shape.text.strip())
            slide_tables = self.extract_tables_from_slide(slide)
            for table in slide_tables:
                slide_content.append(table.to_string(index=False))  # Append table as text
            if slide_content:
                overall_content.append(f"Slide {idx + 1}:\n" + "\n".join(slide_content))
        return "\n\n".join(overall_content)

    def preprocess_ppt(self):
        """
        Extracts and structures text and table data from each slide in the presentation.

        Returns:
            list: A list of dictionaries containing slide title, content, and table data.
        """
        try:
            data = []
            for slide in self.presentation.slides:
                slide_text = []
                slide_tables = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide_text.append(shape.text.strip())
                    elif shape.has_table:
                        table_data = []
                        table = shape.table
                        for row in table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            table_data.append(row_data)
                        slide_tables.append(table_data)
                if slide_text or slide_tables:
                    data.append({
                        "title": slide_text[0] if slide_text else "Untitled Slide",
                        "content": " ".join(slide_text[1:]) if len(slide_text) > 1 else "",
                        "tables": slide_tables
                    })
            return data
        except Exception as e:
            print(f"Error preprocessing presentation: {e}")
            return []
